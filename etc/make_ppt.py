from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

# ── 색상 팔레트 ──────────────────────────────────────────────
NAVY    = RGBColor(0x0D, 0x1B, 0x2A)   # 배경
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW  = RGBColor(0xFF, 0xDA, 0x79)   # 강조
BLUE    = RGBColor(0xA8, 0xD8, 0xEA)   # 서브
GRAY    = RGBColor(0xCC, 0xCC, 0xCC)   # 본문

W = Inches(13.33)   # 와이드 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # 완전 빈 레이아웃


# ── 헬퍼 함수 ────────────────────────────────────────────────
def add_bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_box(slide, left, top, width, height,
            text="", font_size=18, bold=False,
            color=WHITE, align=PP_ALIGN.LEFT,
            bg=None, bg_alpha=None,
            wrap=True, font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = font_name
    if bg:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg
    return txBox

def add_line(slide, left, top, width, color=YELLOW, height=Pt(2)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, int(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def title_slide_layout(slide, title, subtitle=""):
    add_bg(slide)
    add_line(slide,
             left=Inches(1), top=Inches(3.2),
             width=Inches(11.33), color=YELLOW)
    add_box(slide,
            left=Inches(1), top=Inches(1.2),
            width=Inches(11.33), height=Inches(1.8),
            text=title, font_size=40, bold=True,
            color=YELLOW, align=PP_ALIGN.CENTER)
    if subtitle:
        add_box(slide,
                left=Inches(1), top=Inches(3.6),
                width=Inches(11.33), height=Inches(0.8),
                text=subtitle, font_size=20,
                color=GRAY, align=PP_ALIGN.CENTER)

def section_header(slide, text, color=YELLOW):
    add_bg(slide)
    add_box(slide,
            left=Inches(0.5), top=Inches(0.25),
            width=Inches(12.33), height=Inches(0.6),
            text=text, font_size=22, bold=True,
            color=color, align=PP_ALIGN.LEFT)
    add_line(slide,
             left=Inches(0.5), top=Inches(0.9),
             width=Inches(12.33), color=YELLOW)

def row(slide, top, label, value,
        lw=Inches(3.5), vw=Inches(8.8),
        rh=Inches(0.42), fs=14, vc=WHITE, lc=BLUE):
    left_l = Inches(0.5)
    left_v = Inches(4.1)
    add_box(slide, left_l, top, lw, rh,
            text=label, font_size=fs, bold=True, color=lc)
    add_box(slide, left_v, top, vw, rh,
            text=value, font_size=fs, color=vc)


# ════════════════════════════════════════════════════════════
# 슬라이드 1 : 표지
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_line(s, Inches(1), Inches(2.8), Inches(11.33))
add_box(s, Inches(1), Inches(0.6), Inches(11.33), Inches(1.8),
        "스토리보드 기획 문서", font_size=48, bold=True,
        color=YELLOW, align=PP_ALIGN.CENTER)
add_box(s, Inches(1), Inches(2.0), Inches(11.33), Inches(0.7),
        "Storyboard Planning Document", font_size=22,
        color=BLUE, align=PP_ALIGN.CENTER)
add_box(s, Inches(1), Inches(3.1), Inches(11.33), Inches(0.6),
        "캠페인 : 만들기 전에, 한 번 더", font_size=20, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
add_box(s, Inches(1), Inches(3.8), Inches(11.33), Inches(0.5),
        "K-AI 콘텐츠 공모전 2026  |  Track A. 대학(원)/일반부 AI 영상(광고)", font_size=16,
        color=GRAY, align=PP_ALIGN.CENTER)
add_box(s, Inches(1), Inches(6.6), Inches(11.33), Inches(0.5),
        "2026-06-14", font_size=14,
        color=GRAY, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════
# 슬라이드 2 : 브랜드 아이덴티티
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
section_header(s, "01  브랜드 아이덴티티 / Brand Identity")

items = [
    ("캠페인명",         "만들기 전에, 한 번 더"),
    ("공모 주제 분류",   "AI 윤리 준수 — 생성형 AI 사용 시 지켜야 할 에티켓 및 저작권 보호"),
    ("타겟",             "생성형 AI 창작 도구를 사용하는 대학생·크리에이터·일반인 (20~40대)"),
    ("톤앤매너",         "미니멀·감성적 / 차분하고 여운 있는 시네마틱"),
    ("USP (차별점)",     '"AI로 만들 때도 창작자의 권리를 먼저 생각한다"는 태도 자체가 새로운 창작 문화임을 보여줌'),
    ("핵심 메시지",      "AI로 만들 때도, 창작자를 먼저 생각합니다."),
    ("엔딩 슬로건",      "만들기 전에, 한 번 더."),
]
for i, (label, value) in enumerate(items):
    row(s, top=Inches(1.15 + i * 0.72), label=label, value=value)


# ════════════════════════════════════════════════════════════
# 슬라이드 3 : 서사 구조 + 영상 스펙
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
section_header(s, "02  캠페인 목표 및 서사 구조 / Narrative Structure")

add_box(s, Inches(0.5), Inches(1.05), Inches(12.33), Inches(0.45),
        "광고 목적 : AI 저작권·에티켓 인식 제고  (인지 + 태도 변화)", font_size=15, color=GRAY)

narrative = [
    ("기 (씬 1)", "누군가의 손에서 창작물이 탄생한다"),
    ("승 (씬 2)", "그 수많은 창작물들이 AI를 만들었다"),
    ("전 (씬 3)", "이제 당신도 AI로 창작하는 사람이다"),
    ("결 (씬 4·5)", "만들기 전에, 한 번 더 — 출처·허락·존중"),
]
for i, (k, v) in enumerate(narrative):
    t = Inches(1.6 + i * 0.78)
    add_box(s, Inches(0.5), t, Inches(2.2), Inches(0.6),
            k, font_size=16, bold=True, color=YELLOW)
    add_box(s, Inches(2.9), t, Inches(10.0), Inches(0.6),
            v, font_size=16, color=WHITE)

add_line(s, Inches(0.5), Inches(4.8), Inches(12.33), color=BLUE)
add_box(s, Inches(0.5), Inches(4.9), Inches(12.33), Inches(0.45),
        "영상 스펙", font_size=16, bold=True, color=BLUE)
specs = "총 40초  |  1920×1080 (1080p)  |  16:9 가로형  |  24fps  |  H.264 / AAC  |  공모전 공식 지정 음원  |  AI 워터마크 우측 하단"
add_box(s, Inches(0.5), Inches(5.4), Inches(12.33), Inches(0.5),
        specs, font_size=14, color=GRAY)


# ════════════════════════════════════════════════════════════
# 슬라이드 4 : 사용 도구 목록
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
section_header(s, "03  사용 도구 목록 / Tool List")

tools = [
    ("이미지 생성",      "Google Gemini",          "Ideogram / Canva AI",       "프롬프트 접근성이 높고 사실적·일러스트 스타일 모두 지원"),
    ("영상 변환",        "Runway Gen-3",            "Pika",                      "이미지→모션 변환 품질 우수, 씬별 세밀한 모션 제어 가능"),
    ("음성 합성 TTS",    "ElevenLabs",              "CLOVA Voice",               "감성적이고 자연스러운 한국어 목소리 지원"),
    ("BGM",              "공모전 공식 지정 음원",    "—",                         "공모전 규정 필수"),
    ("편집 통합",        "CapCut",                  "DaVinci Resolve",           "텍스트 애니메이션·AI 워터마크 삽입·BGM 페이드아웃"),
]
headers = ["목적", "주 도구", "대안 도구", "선택 이유"]
col_w   = [Inches(1.8), Inches(2.3), Inches(2.3), Inches(5.5)]
col_x   = [Inches(0.4), Inches(2.3), Inches(4.7), Inches(7.1)]

for ci, h in enumerate(headers):
    add_box(s, col_x[ci], Inches(1.05), col_w[ci], Inches(0.45),
            h, font_size=14, bold=True, color=YELLOW)
add_line(s, Inches(0.4), Inches(1.52), Inches(12.53), color=YELLOW, height=Pt(1))

for ri, tool in enumerate(tools):
    t = Inches(1.65 + ri * 0.88)
    colors = [BLUE, WHITE, GRAY, GRAY]
    for ci, val in enumerate(tool):
        add_box(s, col_x[ci], t, col_w[ci], Inches(0.8),
                val, font_size=13, color=colors[ci])


# ════════════════════════════════════════════════════════════
# 씬 슬라이드 생성 함수
# ════════════════════════════════════════════════════════════
def scene_slide(prs, scene_no, time_range, duration,
                objective, composition, narration,
                tools_used, tool_reason,
                prompt_label, prompt_en, prompt_ko,
                runway_prompt, runway_ko,
                filename, result_summary,
                tts_text="", tts_style=""):

    s = prs.slides.add_slide(BLANK)
    section_header(s, f"씬 {scene_no}  ({time_range}, {duration}초)")

    fields_left = [
        ("목표 메시지",   objective),
        ("화면 구성",     composition),
        ("내레이션/카피", narration),
        ("사용 도구",     tools_used),
        ("도구 선택 이유",tool_reason),
    ]
    for i, (lbl, val) in enumerate(fields_left):
        row(s, top=Inches(1.05 + i * 0.72),
            label=lbl, value=val, fs=13)

    # 프롬프트 박스
    add_line(s, Inches(0.5), Inches(4.65), Inches(12.33), color=BLUE)
    add_box(s, Inches(0.5), Inches(4.75), Inches(5.8), Inches(0.35),
            f"이미지 프롬프트 ({prompt_label})", font_size=13, bold=True, color=BLUE)
    add_box(s, Inches(0.5), Inches(5.15), Inches(5.8), Inches(1.9),
            prompt_en, font_size=11, color=WHITE, wrap=True)
    add_box(s, Inches(6.5), Inches(4.75), Inches(6.3), Inches(0.35),
            "한글 요약", font_size=13, bold=True, color=BLUE)
    add_box(s, Inches(6.5), Inches(5.15), Inches(6.3), Inches(1.9),
            prompt_ko, font_size=11, color=GRAY, wrap=True)

    # 파일명 + 결과 요약
    add_box(s, Inches(0.5), Inches(7.1), Inches(6.0), Inches(0.3),
            f"파일명 : {filename}", font_size=11, color=GRAY)
    add_box(s, Inches(6.5), Inches(7.1), Inches(6.3), Inches(0.3),
            f"결과 요약 : {result_summary}", font_size=11, color=GRAY)

    return s


def scene_slide_runway(prs, scene_no, runway_prompt, runway_ko,
                       tts_text="", tts_style=""):
    """Runway + TTS 프롬프트 슬라이드"""
    s = prs.slides.add_slide(BLANK)
    section_header(s, f"씬 {scene_no}  —  영상 변환 & 내레이션 프롬프트")

    add_box(s, Inches(0.5), Inches(1.1), Inches(5.8), Inches(0.4),
            "Runway 영상 변환 프롬프트", font_size=15, bold=True, color=YELLOW)
    add_box(s, Inches(0.5), Inches(1.55), Inches(5.8), Inches(2.2),
            runway_prompt, font_size=13, color=WHITE, wrap=True)
    add_box(s, Inches(6.5), Inches(1.1), Inches(6.3), Inches(0.4),
            "한글 요약", font_size=15, bold=True, color=BLUE)
    add_box(s, Inches(6.5), Inches(1.55), Inches(6.3), Inches(2.2),
            runway_ko, font_size=13, color=GRAY, wrap=True)

    if tts_text:
        add_line(s, Inches(0.5), Inches(3.9), Inches(12.33), color=BLUE)
        add_box(s, Inches(0.5), Inches(4.0), Inches(12.33), Inches(0.4),
                "ElevenLabs 내레이션 텍스트", font_size=15, bold=True, color=BLUE)
        add_box(s, Inches(0.5), Inches(4.5), Inches(12.33), Inches(0.7),
                tts_text, font_size=18, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
        add_box(s, Inches(0.5), Inches(5.3), Inches(12.33), Inches(0.4),
                f"목소리 설정 : {tts_style}", font_size=13, color=GRAY)
    return s


# ════════════════════════════════════════════════════════════
# 슬라이드 5~6 : 씬 1 (이미지 + Runway)
# ════════════════════════════════════════════════════════════
scene_slide(prs,
    scene_no="1  —  창작자의 손 / Creator's Hand  [컷 A·B·C]",
    time_range="0~7초", duration=7,
    objective="사람의 손에서 창작물이 탄생한다는 사실을 시각적으로 각인시킨다",
    composition="클로즈업 / 붓·건반·키보드 손 3컷 / 어두운 스튜디오 배경 / 텍스트 없음",
    narration="없음 (음악과 영상만)",
    tools_used="이미지 : Google Gemini  |  영상 : Runway Gen-3",
    tool_reason="3개 장면(화가·피아니스트·작가) 키비주얼 생성 후 Runway로 슬로우모션 변환",
    prompt_label="Gemini / 컷 A — 화가의 손",
    prompt_en=(
        "A dramatic close-up of a painter's hand holding a fine brush,\n"
        "making a single delicate stroke on a white canvas.\n"
        "Chiaroscuro lighting, dark studio background, warm golden highlights,\n"
        "shallow depth of field, cinematic, photorealistic, 4K."
    ),
    prompt_ko=(
        "화가의 손이 붓으로 한 획을 긋는 클로즈업.\n"
        "명암 대비 강한 조명, 어두운 배경, 황금빛 하이라이트.\n\n"
        "[컷 B] 피아니스트 손, 청백 색조\n"
        "[컷 C] 작가 손, 기계식 키보드 타이핑"
    ),
    runway_prompt=(
        "[컷A] The hand slowly raises the brush and makes a single deliberate stroke,\n"
        "gentle motion, cinematic slow motion, 24fps.\n\n"
        "[컷B] Fingers gently press the piano keys one by one,\n"
        "slow and deliberate, sound wave ripple effect on keys.\n\n"
        "[컷C] Fingers type rhythmically, subtle motion blur on fingers."
    ),
    runway_ko="각 컷 : 슬로우모션으로 손의 창작 행위를 강조",
    filename="scene01A_painter.png / scene01B_pianist.png / scene01C_writer.png\nscene01A_motion.mp4 / scene01B_motion.mp4 / scene01C_motion.mp4",
    result_summary="황금빛·청백·무채색 3가지 손 클로즈업, 창작 행위의 다양성 시각화"
)


# ════════════════════════════════════════════════════════════
# 슬라이드 7~8 : 씬 2
# ════════════════════════════════════════════════════════════
scene_slide(prs,
    scene_no="2  —  AI 학습 시각화 / AI Learning Visualization",
    time_range="7~17초", duration=10,
    objective="수많은 창작물이 AI를 학습시켰다는 사실을 추상적으로 시각화한다",
    composition="광각 / 창작물 파티클 → 중앙 AI 아이콘 수렴 / 어두운 우주 배경 / 내레이션 텍스트",
    narration="\"이 모든 창작물이... AI를 만들었습니다.\"",
    tools_used="이미지 : Gemini  |  영상 : Runway  |  음성 : ElevenLabs",
    tool_reason="추상 파티클 시각화 → Runway 수렴 모션 → ElevenLabs 감성 내레이션",
    prompt_label="Gemini",
    prompt_en=(
        "Abstract digital visualization of thousands of tiny artworks,\n"
        "music notes, handwritten text, and painting fragments\n"
        "all flowing and converging into a single glowing AI neural network icon.\n"
        "Deep dark navy space background, electric blue and white luminous particles,\n"
        "flowing data streams, cinematic, high contrast, dramatic, 4K."
    ),
    prompt_ko=(
        "그림·악보·텍스트 조각들이 중앙의 빛나는 AI 아이콘으로 수렴.\n"
        "어두운 우주 배경, 전기빛 파란색·흰색 발광 파티클, 데이터 스트림."
    ),
    runway_prompt=(
        "Thousands of tiny image and text particles swirl and stream inward,\n"
        "converging into the central glowing AI icon which pulses with light,\n"
        "smooth flowing motion, cinematic VFX style."
    ),
    runway_ko="파티클이 소용돌이치며 AI 아이콘으로 수렴, 아이콘이 맥박처럼 빛남",
    filename="scene02_ai_learning.png  /  scene02_motion.mp4  /  scene02_narration.mp3",
    result_summary="어두운 우주 배경 위 파란 파티클 수렴 이미지 확보, 광고 톤앤매너 일치",
    tts_text="이 모든 창작물이... AI를 만들었습니다.",
    tts_style="차분하고 낮은 톤 / 느린 속도 / 감성적"
)
scene_slide_runway(prs, scene_no=2,
    runway_prompt=(
        "Thousands of tiny image and text particles swirl and stream inward,\n"
        "converging into the central glowing AI icon which pulses with light,\n"
        "smooth flowing motion, cinematic VFX style."
    ),
    runway_ko="파티클이 소용돌이치며 AI 아이콘으로 수렴, 아이콘이 맥박처럼 빛남",
    tts_text="이 모든 창작물이...  AI를 만들었습니다.",
    tts_style="차분하고 낮은 톤 / 느린 속도 / 감성적"
)


# ════════════════════════════════════════════════════════════
# 슬라이드 9~10 : 씬 3
# ════════════════════════════════════════════════════════════
scene_slide(prs,
    scene_no="3  —  당신도 창작자 / You Are Also a Creator",
    time_range="17~25초", duration=8,
    objective="AI 도구를 사용하는 사람도 창작 행위의 주체임을 인식시킨다",
    composition="미디엄 샷 / 노트북 앞 인물 실루엣 / 따뜻한 창문 자연광 / 내레이션 텍스트",
    narration="\"이제, 당신도 창작자입니다.\"",
    tools_used="이미지 : Gemini  |  영상 : Runway  |  음성 : ElevenLabs",
    tool_reason="실루엣 처리로 시청자가 인물에 자신을 투영, 자연광으로 씬2 어두움 → 희망으로 전환",
    prompt_label="Gemini",
    prompt_en=(
        "A medium shot of a person's silhouette sitting at a desk with a laptop,\n"
        "the laptop screen glowing brightly showing colorful AI-generated artwork.\n"
        "Soft warm sunlight streaming through a window behind the figure,\n"
        "golden dust particles floating in the light, lens flare,\n"
        "cinematic, hopeful atmosphere, photorealistic, 4K."
    ),
    prompt_ko=(
        "노트북 앞에 앉은 인물 실루엣.\n"
        "창문에서 흘러드는 따뜻한 자연광, 먼지 파티클, 렌즈 플레어.\n"
        "희망적이고 고양된 분위기."
    ),
    runway_prompt=(
        "The person slowly raises their head and looks at the glowing screen,\n"
        "camera slowly pushes in toward the laptop screen,\n"
        "warm sunlight gradually brightens, cinematic slow motion."
    ),
    runway_ko="인물이 고개를 들어 화면을 바라보며 카메라가 화면 쪽으로 천천히 당겨짐",
    filename="scene03_creator.png  /  scene03_motion.mp4  /  scene03_narration.mp3",
    result_summary="따뜻한 자연광 아래 노트북 앞 실루엣, 희망적 분위기 전환 확보"
)
scene_slide_runway(prs, scene_no=3,
    runway_prompt=(
        "The person slowly raises their head and looks at the glowing screen,\n"
        "camera slowly pushes in toward the laptop screen,\n"
        "warm sunlight gradually brightens, cinematic slow motion."
    ),
    runway_ko="인물이 고개를 들어 화면을 바라보며 카메라가 화면 쪽으로 천천히 당겨짐",
    tts_text="이제,  당신도 창작자입니다.",
    tts_style="따뜻하고 부드러운 톤 / 확신 있는 속도"
)


# ════════════════════════════════════════════════════════════
# 슬라이드 11~12 : 씬 4
# ════════════════════════════════════════════════════════════
scene_slide(prs,
    scene_no="4  —  행동 지침 / Three Actions",
    time_range="25~35초", duration=10,
    objective="출처 표기·허락·존중, 3가지 구체적 행동을 명확하게 전달한다",
    composition="센터 프레임 / 아티스트·AI사용자 미니멀 일러스트 / 흰 배경 / 텍스트 3단계 순차 강조",
    narration="\"출처를 밝히고, 허락을 구하고, 창작자를 존중합니다.\"",
    tools_used="이미지 : Gemini  |  영상 : Runway  |  음성 : ElevenLabs  |  편집 : CapCut",
    tool_reason="Gemini 미니멀 일러스트 생성, CapCut 3단계 텍스트 페이드인 삽입",
    prompt_label="Gemini",
    prompt_en=(
        "A minimalist flat vector illustration of two friendly figures standing side by side.\n"
        "Left: an artist holding a paintbrush and color palette.\n"
        "Right: a person using a laptop with a small glowing star icon on the screen.\n"
        "Clean white background, pastel soft blue and warm yellow,\n"
        "simple geometric shapes, balanced composition, no text, no shadows."
    ),
    prompt_ko=(
        "두 인물이 나란히 선 미니멀 플랫 일러스트.\n"
        "왼쪽: 붓과 팔레트를 든 아티스트.\n"
        "오른쪽: 노트북 AI 사용자.\n"
        "흰 배경, 파스텔 블루·노란색, 텍스트·그림자 없음."
    ),
    runway_prompt=(
        "The two figures simultaneously turn toward each other and nod gently,\n"
        "subtle friendly animation, soft pastel glow between them,\n"
        "smooth looping motion, flat 2D animation style."
    ),
    runway_ko="두 인물이 서로를 향해 고개를 끄덕이는 친근한 애니메이션, 사이에 파스텔 빛",
    filename="scene04_coexist.png  /  scene04_motion.mp4  /  scene04_narration.mp3",
    result_summary="파스텔 톤 두 인물 공존 일러스트, 따뜻하고 균형 있는 구성 확보"
)
scene_slide_runway(prs, scene_no=4,
    runway_prompt=(
        "The two figures simultaneously turn toward each other and nod gently,\n"
        "subtle friendly animation, soft pastel glow between them,\n"
        "smooth looping motion, flat 2D animation style."
    ),
    runway_ko="두 인물이 서로를 향해 고개를 끄덕이는 애니메이션",
    tts_text="출처를 밝히고,  허락을 구하고,  창작자를 존중합니다.",
    tts_style="명확하고 또렷한 톤 / 세 문장 사이 짧은 간격"
)


# ════════════════════════════════════════════════════════════
# 슬라이드 13 : 씬 5 — 엔딩 카드
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
section_header(s, "씬 5  —  엔딩 슬로건 카드 / Ending Card  (35~40초, 5초)")

fields = [
    ("씬 번호",       "5"),
    ("씬 길이",       "5초  (35~40초)"),
    ("목표 메시지",   "슬로건을 각인시키고 AI 워터마크와 함께 공모전 아이덴티티를 마무리한다"),
    ("사용 도구",     "편집 : CapCut  |  음성 : ElevenLabs"),
    ("AI 이미지 생성","해당 없음 (편집 작업)"),
]
for i, (lbl, val) in enumerate(fields):
    row(s, top=Inches(1.05 + i * 0.68), label=lbl, value=val, fs=13)

add_line(s, Inches(0.5), Inches(4.6), Inches(12.33), color=BLUE)
add_box(s, Inches(0.5), Inches(4.7), Inches(12.33), Inches(0.4),
        "CapCut 설정", font_size=15, bold=True, color=BLUE)
capcut_info = (
    "배경색 #0D1B2A (네이비)  |  메인 텍스트 : 만들기 전에, 한 번 더.  |  색상 #FFFFFF  |  폰트 Noto Sans KR Bold\n"
    "서브 텍스트 : AI로 만들 때도, 창작자를 먼저 생각합니다.  |  등장 효과 : 페이드인 0.5초\n"
    "BGM : 공모전 공식 지정 음원 페이드아웃  |  AI 워터마크 : 우측 하단 삽입"
)
add_box(s, Inches(0.5), Inches(5.2), Inches(12.33), Inches(1.0),
        capcut_info, font_size=13, color=GRAY, wrap=True)

add_box(s, Inches(0.5), Inches(6.4), Inches(5.0), Inches(0.4),
        "ElevenLabs 입력 텍스트", font_size=13, bold=True, color=BLUE)
add_box(s, Inches(0.5), Inches(6.85), Inches(5.0), Inches(0.45),
        "만들기 전에, 한 번 더.", font_size=16, bold=True, color=YELLOW)
add_box(s, Inches(6.5), Inches(6.4), Inches(6.3), Inches(0.45),
        "목소리 설정 : 여운이 있는 낮고 조용한 톤 / 천천히", font_size=13, color=GRAY)


# ════════════════════════════════════════════════════════════
# 슬라이드 14 : 프롬프트 수정 전/후 기록
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
section_header(s, "04  프롬프트 수정 전/후 기록 / Prompt Revision Log  (씬 2 기준)")

log = [
    ("수정 전 프롬프트", "many images and text flowing into an AI icon"),
    ("수정 전 문제점",   "단순한 카툰 스타일로 생성됨 — 광고 톤앤매너와 불일치, 파티클 규모감·긴장감 부족"),
    ("수정 방향",        "① 시네마틱 키워드 추가  ② 배경 색상 명시  ③ 파티클 종류·발광 속성 구체화  ④ 수렴 행위를 동사로 명확히 표현"),
    ("수정 후 프롬프트", "Abstract digital visualization of thousands of tiny artworks, music notes, handwritten text, and painting fragments all flowing and converging into a single glowing AI neural network icon. Deep dark navy space background, electric blue and white luminous particles, cinematic, 4K."),
    ("수정 결과",        "어두운 우주 배경 위 파란 발광 파티클이 AI 아이콘으로 수렴하는 이미지 확보. 광고 톤앤매너와 일치."),
]
for i, (lbl, val) in enumerate(log):
    c = YELLOW if "수정 후" in lbl or "결과" in lbl else BLUE
    row(s, top=Inches(1.05 + i * 0.95), label=lbl, value=val, fs=13,
        rh=Inches(0.82), vc=WHITE, lc=c)


# ════════════════════════════════════════════════════════════
# 슬라이드 15 : 타임라인 + 체크리스트
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
section_header(s, "05  타임라인 요약 / Timeline & 제작 체크리스트")

timeline = [
    ("씬 1", "0~7초",   "7초",  "창작자의 손 (화가·피아니스트·작가 3컷)",       "Gemini + Runway"),
    ("씬 2", "7~17초",  "10초", "창작물 → AI 수렴 시각화 + 내레이션",           "Gemini + Runway + ElevenLabs"),
    ("씬 3", "17~25초", "8초",  "실루엣 창작자 + 분위기 전환 + 내레이션",        "Gemini + Runway + ElevenLabs"),
    ("씬 4", "25~35초", "10초", "출처·허락·존중 3단계 행동 지침",                "Gemini + Runway + ElevenLabs + CapCut"),
    ("씬 5", "35~40초", "5초",  "슬로건 엔딩 카드 + AI 워터마크",                "ElevenLabs + CapCut"),
]
headers = ["씬", "구간", "길이", "핵심 내용", "주요 도구"]
col_x2  = [Inches(0.4), Inches(1.3), Inches(2.3), Inches(3.2), Inches(8.5)]
col_w2  = [Inches(0.8), Inches(0.9), Inches(0.8), Inches(5.2), Inches(3.9)]

for ci, h in enumerate(headers):
    add_box(s, col_x2[ci], Inches(1.05), col_w2[ci], Inches(0.4),
            h, font_size=13, bold=True, color=YELLOW)
add_line(s, Inches(0.4), Inches(1.47), Inches(12.53), color=YELLOW, height=Pt(1))

for ri, tl in enumerate(timeline):
    t = Inches(1.6 + ri * 0.68)
    cs = [BLUE, GRAY, GRAY, WHITE, GRAY]
    for ci, val in enumerate(tl):
        add_box(s, col_x2[ci], t, col_w2[ci], Inches(0.6),
                val, font_size=12, color=cs[ci])

add_box(s, Inches(0.4), Inches(5.35), Inches(12.53), Inches(0.4),
        "합계 : 40초  |  1920×1080 / 16:9 / H.264 / AAC  |  공모전 공식 지정 음원  |  AI 워터마크 필수",
        font_size=13, bold=True, color=YELLOW)

add_line(s, Inches(0.4), Inches(5.85), Inches(12.53), color=BLUE)
add_box(s, Inches(0.4), Inches(5.95), Inches(12.53), Inches(0.4),
        "출품 일정 : 공모 기간 2026.06.11~07.15  |  접수 기간 2026.07.01~07.15  |  SNS 업로드 후 URL 제출",
        font_size=13, color=GRAY)


# ════════════════════════════════════════════════════════════
# 저장
# ════════════════════════════════════════════════════════════
OUTPUT = r"c:\Users\YSW1\Desktop\20260601 codyssey_mission\N_B1-2_Contents\01_document\storyboard.pptx"
prs.save(OUTPUT)
print(f"저장 완료 : {OUTPUT}")
